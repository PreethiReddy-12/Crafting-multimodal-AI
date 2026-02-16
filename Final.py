# app.py
# GenAI Craft Fusion — Unified Multi-Feature Streamlit App
# Run with: streamlit run app.py

import os
import json
import time
import requests
import uuid
import whisper
import streamlit as st
import pdfplumber
import docx
from pptx import Presentation
from transformers import pipeline, BlipProcessor, BlipForConditionalGeneration
from PIL import Image
import torch
from moviepy.editor import AudioFileClip, VideoFileClip
import yt_dlp

# -----------------------------
# Global App Config
# -----------------------------
st.set_page_config(page_title="GenAI Craft Fusion", page_icon="🧪", layout="wide")
st.title("GenAI Craft Fusion — Unified App")

tabs = st.tabs([
    "💬 Text Chat",
    "📄 Document Analyzer",
    "🖼 Image Captioning",
    "🎧 Audio Player",
    "🎬 Video Viewer"
])

# -----------------------------
# Tab 1: Text Chat (Ollama)
# -----------------------------
with tabs[0]:
    SYSTEM_PROMPT = (
        "You are Copilot in GenAI Craft Fusion, answering clearly and concisely with streaming output. "
        "Focus on correctness, minimal redundancy, and practical value. If something is ambiguous, ask a brief clarifying question."
    )

    def stream_ollama(prompt, model):
        url = "http://127.0.0.1:11434/api/generate"
        payload = {
            "model": model,
            "prompt": prompt,
            "stream": True,
            "options": {"temperature": 0.2}
        }
        try:
            with requests.post(url, json=payload, stream=True, timeout=30) as r:
                r.raise_for_status()
                for line in r.iter_lines(decode_unicode=True):
                    if not line:
                        continue
                    try:
                        data = json.loads(line)
                    except json.JSONDecodeError:
                        continue
                    if "response" in data and data["response"]:
                        yield data["response"]
                    if data.get("done"):
                        break
        except requests.exceptions.RequestException as e:
            raise RuntimeError(f"Ollama streaming error: {e}")

    st.caption("Real-time answers using Ollama locally. No API key required.")
    default_model = "llama3.1"
    model = st.text_input("Ollama model", default_model)
    stream_delay_ms = st.slider("Output pacing (ms per chunk)", 0, 100, 0)

    if "chat" not in st.session_state:
        st.session_state.chat = [{"role": "system", "content": SYSTEM_PROMPT}]

    user_input = st.text_area("Ask anything:", placeholder="e.g., Explain binary search with a Python example.")
    col1, col2 = st.columns([1, 1])
    with col1:
        send = st.button("Ask")
    with col2:
        clear = st.button("Clear")

    if clear:
        st.session_state.chat = [{"role": "system", "content": SYSTEM_PROMPT}]
        st.experimental_rerun()

    for msg in st.session_state.chat:
        if msg["role"] == "system":
            continue
        with st.chat_message("assistant" if msg["role"] == "assistant" else "user"):
            st.write(msg["content"])

    if send and user_input.strip():
        st.session_state.chat.append({"role": "user", "content": user_input})
        messages = st.session_state.chat
        placeholder = st.empty()
        streamed_text = ""

        try:
            history = []
            for m in messages:
                if m["role"] == "system":
                    history.append(f"System: {m['content']}")
                elif m["role"] == "user":
                    history.append(f"User: {m['content']}")
                else:
                    history.append(f"Assistant: {m['content']}")
            prompt = "\n".join(history) + "\nAssistant:"

            token_stream = stream_ollama(prompt, model=model)

            with st.chat_message("assistant"):
                for token in token_stream:
                    streamed_text += token
                    placeholder.markdown(streamed_text)
                    if stream_delay_ms > 0:
                        time.sleep(stream_delay_ms / 1000.0)

            st.session_state.chat.append({"role": "assistant", "content": streamed_text})

        except Exception as e:
            with st.chat_message("assistant"):
                st.error(f"Streaming failed: {e}")

# -----------------------------
# Tab 2: Document Analyzer
# -----------------------------
with tabs[1]:
    st.caption("Upload PDF, DOCX, or PPTX and click Analyze.")
    @st.cache_resource
    def load_summarizer():
        return pipeline("summarization", model="facebook/bart-large-cnn")
    summarizer = load_summarizer()

    def extract_pdf(uploaded_file):
        text = ""
        with pdfplumber.open(uploaded_file) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text.strip()

    def extract_docx(uploaded_file):
        doc = docx.Document(uploaded_file)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text.strip()

    def extract_pptx(uploaded_file):
        prs = Presentation(uploaded_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()

    uploaded_file = st.file_uploader("Upload a document", type=["pdf", "docx", "pptx"])
    user_prompt = st.text_area("Describe what you want from the document:")

    if st.button("Analyze") and uploaded_file is not None:
        file_type = uploaded_file.name.split(".")[-1].lower()
        if file_type == "pdf":
            doc_text = extract_pdf(uploaded_file)
        elif file_type == "docx":
            doc_text = extract_docx(uploaded_file)
        elif file_type == "pptx":
            doc_text = extract_pptx(uploaded_file)
        else:
            st.error("Unsupported file type.")
            doc_text = ""

        if doc_text:
            st.subheader("Extracted Text (first 1000 chars)")
            st.text(doc_text[:1000])
            analysis_text = f"{user_prompt}\n\nDocument:\n{doc_text}"
            max_chunk = 1000
            chunks = [analysis_text[i:i+max_chunk] for i in range(0, len(analysis_text), max_chunk)]
            summary = ""
            for chunk in chunks:
                if len(chunk.strip()) > 0:
                    result = summarizer(chunk, max_length=200, min_length=50, do_sample=False)
                    summary += result[0]['summary_text'] + "\n"
            st.subheader("Analysis Result")
            st.write(summary.strip())

# -----------------------------
# Tab 3: Image Captioning
# -----------------------------
with tabs[2]:
    uploaded_file = st.file_uploader("Upload an image", type=["jpg", "jpeg", "png"])
    user_prompt = st.text_area("Describe what you want from the image:")

    if st.button("Analyze Image") and uploaded_file is not None:
        image = Image.open(uploaded_file).convert("RGB")

        @st.cache_resource
        def load_model():
            processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base", use_fast=True)
            model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
            return processor, model

        processor, model = load_model()
        inputs = processor(images=image, return_tensors="pt")
        out = model.generate(**inputs)
        caption = processor.decode(out[0], skip_special_tokens=True)

        st.write("### Image Description")
        st.write(f"**User Intent:** {user_prompt}")
        st.write(f"**Caption:** {caption}")
# -----------------------------
# Tab 4: Audio Player
# -----------------------------
# -----------------------------
# Tab 4: Audio Player
# -----------------------------
# -----------------------------
# Tab 4: Audio Player
# -----------------------------
with tabs[3]:
    audio_url = st.text_input("Enter local audio file path or YouTube URL:")
    user_prompt = st.text_area("Describe what you want from the audio:")

    def download_audio(url):
        out_file = "downloaded_audio.%(ext)s"
        ydl_opts = {
            "format": "bestaudio/best",
            "outtmpl": out_file,
            "postprocessors": [{
                "key": "FFmpegExtractAudio",
                "preferredcodec": "mp3",
                "preferredquality": "192",
            }],
        }
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
        return "downloaded_audio.mp3", info

    if st.button("Analyze Audio") and audio_url:
        try:
            if os.path.exists(audio_url):
                path = audio_url
                info = {"title": os.path.basename(path), "uploader": "Local File", "description": ""}
            else:
                st.info("Downloading audio...")
                path, info = download_audio(audio_url)

            clip = AudioFileClip(path)
            max_duration = 600
            if clip.duration > max_duration:
                clip = clip.subclip(0, max_duration)

            output_file = "output_10min_audio.mp3"
            clip.write_audiofile(output_file)
            st.success("Processed audio ready!")
            st.audio(output_file)

            # --- NEW: Local Whisper transcription ---
            import whisper
            model = whisper.load_model("base")  # "small", "medium", "large" available
            result = model.transcribe(output_file)
            transcript = result["text"]

            # Summarize transcript locally (using HuggingFace model you already use)
            from transformers import pipeline
            summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
            summary = summarizer(transcript, max_length=200, min_length=50, do_sample=False)[0]['summary_text']

            st.write("### Audio Analysis Result")
            st.write(f"**User Intent:** {user_prompt}")
            st.write(f"*Title:* {info.get('title')}")
            st.write(f"*Channel:* {info.get('uploader')}")
            st.write(f"*Description:*\n{info.get('description')}")
            st.write("### Transcript (first 500 chars)")
            st.text(transcript[:500])
            st.write("### Summary of Audio Content")
            st.write(summary)

            clip.close()
        except Exception as e:
            st.error(f"Error: {e}")
# -----------------------------
# Tab 5: Video Viewer
# -----------------------------
# -----------------------------
with tabs[4]:
    video_url = st.text_input("Enter local video path or YouTube URL:")
    user_prompt = st.text_area("Describe what you want from the video:")


    def download_video(url):
        out_file = f"video_{uuid.uuid4().hex}.mp4"
        ydl_opts = {"outtmpl": out_file, "format": "best"}
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])
        return out_file

    if st.button("Analyze Video") and video_url:
        try:
            # --- Download or load local video ---
            if os.path.exists(video_url):
                path = video_url
            else:
                st.info("Downloading video...")
                path = download_video(video_url)

            clip = VideoFileClip(path)

            # --- Cap duration at 10 minutes ---
            max_duration = 600
            if clip.duration > max_duration:
                clip = clip.subclip(0, max_duration)

            # --- Save processed video ---
            output_file = f"processed_{uuid.uuid4().hex}.mp4"
            clip.write_videofile(output_file, codec="libx264", audio_codec="aac", fps=30)

            st.success("Processed video ready!")
            st.video(output_file)

            # --- Extract audio ---
            audio_output = f"audio_{uuid.uuid4().hex}.mp3"
            clip.audio.write_audiofile(audio_output)
            clip.close()  # release file handle

            # --- Transcribe with Whisper ---
            model = whisper.load_model("base")
            result = model.transcribe(audio_output)
            transcript = result.get("text", "").strip()

            if not transcript:
                st.error("No transcript generated — check if the video has audio.")
            else:
                # --- Summarize transcript ---
                summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

                # Chunk transcript if very long
                def chunk_text(text, max_words=500):
                    words = text.split()
                    for i in range(0, len(words), max_words):
                        yield " ".join(words[i:i+max_words])

                summaries = []
                for chunk in chunk_text(transcript):
                    s = summarizer(chunk, max_length=200, min_length=50, do_sample=False)[0]['summary_text']
                    summaries.append(s)

                final_summary = " ".join(summaries)

                # --- Display results ---
                st.write("### Video Analysis Result")
                st.write(f"**User Intent:** {user_prompt}")
                st.write(f"- Duration: {clip.duration:.2f} seconds (capped at 10 minutes)")
                st.write(f"- Resolution: {clip.w} x {clip.h}")
                st.write(f"- FPS: {clip.fps}")
                st.write("### Transcript (first 500 chars)")
                st.text(transcript[:500])
                st.write("### Summary of Video Content")
                st.write(final_summary)

        except Exception as e:
            st.error(f"Error: {e}")